#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Generate the Korean 'IDE for Omics Researcher' (오믹스 연구자를 위한 AI IDE) pptx deck."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---- palette ----
NAVY   = RGBColor(0x0F, 0x2A, 0x43)
NAVY2  = RGBColor(0x16, 0x3B, 0x5C)
TEAL   = RGBColor(0x0F, 0xB5, 0xA9)
GREEN  = RGBColor(0x36, 0xB3, 0x7E)
AMBER  = RGBColor(0xF2, 0xA6, 0x2E)
CORAL  = RGBColor(0xE8, 0x6A, 0x5C)
LIGHT  = RGBColor(0xF4, 0xF7, 0xFA)
CARD   = RGBColor(0xFF, 0xFF, 0xFF)
INK    = RGBColor(0x23, 0x33, 0x42)
GRAY   = RGBColor(0x5A, 0x6B, 0x7B)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LINE   = RGBColor(0xD8, 0xE0, 0xE8)

FONT = "Malgun Gothic"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

def slide():
    return prs.slides.add_slide(BLANK)

def _set_font(run, size, color, bold, italic=False, font=FONT):
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        e = rPr.find(qn(tag))
        if e is None:
            e = rPr.makeelement(qn(tag), {})
            rPr.append(e)
        e.set("typeface", font)

def rect(s, x, y, w, h, fill, line=None, line_w=None, shape=MSO_SHAPE.RECTANGLE):
    sp = s.shapes.add_shape(shape, x, y, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w or 1)
    sp.shadow.inherit = False
    return sp

def txt(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        space_after=6, line_spacing=1.05, wrap=True):
    """runs: list of paragraphs; each paragraph is list of (text,size,color,bold,italic)."""
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0
    tf.margin_top = 0; tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        p.line_spacing = line_spacing
        for (t, sz, col, bold, *rest) in para:
            it = rest[0] if rest else False
            r = p.add_run(); r.text = t
            _set_font(r, sz, col, bold, it)
    return tb

def bg(s, color=LIGHT):
    rect(s, 0, 0, SW, SH, color)

def kicker_bar(s):
    rect(s, 0, 0, SW, Inches(0.14), TEAL)

def header(s, num, title, kicker=None):
    """standard content-slide header."""
    kicker_bar(s)
    # slide number chip
    rect(s, Inches(0.55), Inches(0.5), Inches(0.62), Inches(0.62), NAVY,
         shape=MSO_SHAPE.OVAL)
    txt(s, Inches(0.55), Inches(0.5), Inches(0.62), Inches(0.62),
        [[(f"{num:02d}", 20, WHITE, True)]], align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE)
    paras = []
    if kicker:
        paras.append([(kicker, 12, TEAL, True)])
    paras.append([(title, 26, NAVY, True)])
    txt(s, Inches(1.4), Inches(0.42), Inches(11.3), Inches(0.95), paras,
        space_after=1, line_spacing=1.0, anchor=MSO_ANCHOR.MIDDLE)
    # underline
    rect(s, Inches(1.42), Inches(1.32), Inches(2.2), Pt(3), TEAL)

def footer(s):
    txt(s, Inches(0.55), Inches(7.02), Inches(9), Inches(0.35),
        [[("오믹스 연구자를 위한 AI IDE  ·  IDE for Omics Researcher", 9, GRAY, False)]])
    txt(s, Inches(10.3), Inches(7.02), Inches(2.5), Inches(0.35),
        [[("GHBio Co-Scientist", 9, GRAY, True)]], align=PP_ALIGN.RIGHT)

def bullets(s, x, y, w, items, size=15, gap=10, color=INK, marker="—", mcolor=TEAL):
    tb = s.shapes.add_textbox(x, y, w, Inches(4.5))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left=0; tf.margin_top=0; tf.margin_right=0
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.space_after = Pt(gap); p.line_spacing = 1.08
        if isinstance(it, tuple):
            lead, rest = it
        else:
            lead, rest = None, it
        rm = p.add_run(); rm.text = f"{marker}  "
        _set_font(rm, size, mcolor, True)
        if lead:
            r1 = p.add_run(); r1.text = lead + "  "
            _set_font(r1, size, NAVY, True)
        r2 = p.add_run(); r2.text = rest
        _set_font(r2, size, color, False)
    return tb

def flow(s, steps, y=Inches(2.7), h=Inches(1.7), colors=None):
    """Horizontal box→arrow→box chain. steps: list of (line1,line2)."""
    n = len(steps)
    margin = Inches(0.7)
    arrow_w = Inches(0.55)
    total = SW - 2*margin
    box_w = (total - (n-1)*arrow_w) / n
    x = margin
    palette = colors or [NAVY, NAVY2, TEAL, GREEN, AMBER, CORAL, NAVY, TEAL]
    for i,(l1,l2) in enumerate(steps):
        col = palette[i % len(palette)]
        rect(s, x, y, box_w, h, col, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        paras = [[(l1, 15, WHITE, True)]]
        if l2:
            paras.append([(l2, 11, RGBColor(0xDF,0xEC,0xF0), False)])
        txt(s, x+Inches(0.12), y, box_w-Inches(0.24), h, paras,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=3,
            line_spacing=1.0)
        if i < n-1:
            ax = x + box_w + Emu(int(arrow_w*0.08))
            ar = rect(s, x+box_w, y+h/2-Inches(0.18), arrow_w, Inches(0.36),
                      AMBER, shape=MSO_SHAPE.RIGHT_ARROW)
        x += box_w + arrow_w

def note(s, text, y=Inches(4.9), color=NAVY):
    box = rect(s, Inches(0.7), y, SW-Inches(1.4), Inches(1.0), CARD,
               line=LINE, line_w=1, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s, Inches(0.7), y, Inches(0.12), Inches(1.0), color)
    txt(s, Inches(1.0), y, SW-Inches(2.0), Inches(1.0),
        [[("▸  ", 14, color, True),(text, 14, INK, False)]],
        anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.1)

def cards(s, items, y=Inches(2.0), h=Inches(4.2), cols=None, top_colors=None):
    """Grid of cards. items: list of (title, body)."""
    n = len(items); cols = cols or n
    rows = (n + cols - 1)//cols
    margin = Inches(0.6); gap = Inches(0.35)
    cw = (SW - 2*margin - (cols-1)*gap)/cols
    ch = (h - (rows-1)*gap)/rows
    pal = top_colors or [TEAL, GREEN, AMBER, CORAL, NAVY2, TEAL, GREEN, AMBER, CORAL]
    for idx,(t,b) in enumerate(items):
        r = idx//cols; c = idx%cols
        x = margin + c*(cw+gap); yy = y + r*(ch+gap)
        rect(s, x, yy, cw, ch, CARD, line=LINE, line_w=1,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        rect(s, x, yy, cw, Inches(0.12), pal[idx%len(pal)])
        txt(s, x+Inches(0.22), yy+Inches(0.28), cw-Inches(0.44), ch-Inches(0.4),
            [[(t, 15, NAVY, True)], [(b, 12, GRAY, False)]],
            space_after=6, line_spacing=1.08)

# =========================================================
# SLIDE 1 — What is an IDE?
# =========================================================
s = slide(); bg(s); kicker_bar(s)
header(s, 1, "IDE란 무엇인가 — 소프트웨어 개발의 관점", "1990~2000년대  ·  데스크톱 시대")
bullets(s, Inches(0.7), Inches(1.7), Inches(6.0), [
    ("IDE = 통합 개발 환경 (Integrated Development Environment).", "에디터 + 컴파일러 + 디버거 + 빌드 도구를 하나의 프로그램에."),
    ("코드를 '작업대(workbench)'로 바꿔줍니다 —", "작성·실행·디버깅·배포를 하나의 창에서."),
    ("대표주자:  Microsoft Visual Studio", "이 카테고리를 정의한 기준점 IDE."),
], size=15, gap=14)
# right column flow
txt(s, Inches(7.1), Inches(1.75), Inches(5.5), Inches(0.4),
    [[("Visual Studio 스택", 14, TEAL, True)]])
sub = [("Visual Studio","IDE"),("C++ / C#","프로그래밍 언어"),("Windows OS","대상 플랫폼")]
n=len(sub); yv=Inches(2.35)
for i,(a,b) in enumerate(sub):
    yy = yv + Inches(i*1.25)
    rect(s, Inches(7.1), yy, Inches(5.3), Inches(1.0), NAVY if i==0 else NAVY2,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(s, Inches(7.4), yy, Inches(5.0), Inches(1.0),
        [[(a, 18, WHITE, True),("   "+b, 12, RGBColor(0xBF,0xD6,0xE0), False)]],
        anchor=MSO_ANCHOR.MIDDLE)
    if i<n-1:
        rect(s, Inches(9.5), yy+Inches(1.0)-Inches(0.02), Inches(0.36), Inches(0.28),
             AMBER, shape=MSO_SHAPE.DOWN_ARROW)
note(s, "한 시대의 개발은 'IDE + 언어 + OS'의 조합으로 정의되었습니다 — Visual Studio + C++/C# + Windows.", y=Inches(6.05))
footer(s)

# =========================================================
# SLIDE 2 — VS Code -> Python -> ML
# =========================================================
s = slide(); bg(s); header(s, 2, "VS Code의 등장 — 가볍고, 크로스플랫폼", "2015년~  ·  세상을 삼킨 에디터")
flow(s, [("VS Code","무료·확장 가능한 에디터"),
         ("Python","데이터 & 스크립팅"),
         ("머신러닝","NumPy · PyTorch · scikit-learn")],
     y=Inches(2.5), h=Inches(1.9), colors=[NAVY, TEAL, GREEN])
bullets(s, Inches(0.7), Inches(4.7), Inches(11.8), [
    ("무거운 IDE가 아니라 — 확장(extension)으로 강력해지는 빠른 에디터.", ""),
    ("Python + Jupyter 덕분에 VS Code는 데이터 과학·ML의 기본 터전이 되었습니다.", ""),
], size=15, gap=12)
note(s, "가볍고 무료이며 확장(extension)으로 무한히 커지는 구조 — VS Code가 사실상의 표준 에디터가 됩니다.", y=Inches(5.9))
footer(s)

# =========================================================
# SLIDE 3 — Remote SSH -> Ubuntu -> multi-language
# =========================================================
s = slide(); bg(s); header(s, 3, "원격 개발 — 노트북이 서버를 조종한다", "Remote-SSH  ·  하나의 에디터로 어떤 머신이든")
flow(s, [("VS Code","로컬 UI"),
         ("Remote / SSH","보안 터널"),
         ("Ubuntu OS","강력한 리눅스 서버")],
     y=Inches(1.85), h=Inches(1.55), colors=[NAVY, TEAL, GREEN])
txt(s, Inches(0.7), Inches(3.7), Inches(11.8), Inches(0.4),
    [[("하나의 원격 에디터, 세 가지 소프트웨어:", 15, NAVY, True)]])
cards(s, [
    ("C / C++","네이티브 고성능 애플리케이션"),
    ("Python","머신러닝 & 데이터 파이프라인"),
    ("HTML / JavaScript","웹 애플리케이션 & 대시보드"),
], y=Inches(4.15), h=Inches(2.2), cols=3, top_colors=[CORAL, GREEN, AMBER])
footer(s)

# =========================================================
# SLIDE 4 — VS Code + Copilot -> AI development
# =========================================================
s = slide(); bg(s); header(s, 4, "VS Code + Copilot — 최초의 AI 페어 프로그래머", "2021년~  ·  AI가 에디터 안으로")
flow(s, [("VS Code","에디터"),
         ("GitHub Copilot","LLM 자동완성"),
         ("AI 보조 개발","타이핑하는 대로 코드 제안")],
     y=Inches(2.5), h=Inches(1.9), colors=[NAVY, TEAL, AMBER])
bullets(s, Inches(0.7), Inches(4.7), Inches(11.8), [
    ("대규모 언어모델(LLM)이 코드를 인라인으로 완성 — 대중화된 최초의 '에디터 속 AI'.", ""),
    ("개발자는 반복적인 상용구를 직접 치지 않습니다. 모델이 초안을, 사람은 검토를.", ""),
], size=15, gap=12)
note(s, "Copilot은 'AI가 코드를 제안하는' 시대의 문을 열었습니다 — 개발의 무게중심이 타이핑에서 검토로 이동.", y=Inches(5.9))
footer(s)

# =========================================================
# SLIDE 5 — VS Code + multi LLM -> Cursor
# =========================================================
s = slide(); bg(s); header(s, 5, "Cursor — 여러 LLM을 중심으로 다시 태어난 VS Code", "2023년~  ·  Anysphere의 AI 네이티브 포크")
flow(s, [("VS Code 기반","오픈소스 코어"),
         ("멀티 LLM","Claude · GPT · Gemini"),
         ("Cursor (Anysphere)","AI 네이티브 IDE")],
     y=Inches(1.65), h=Inches(1.35), colors=[NAVY, TEAL, GREEN])
txt(s, Inches(0.7), Inches(3.2), Inches(11.8), Inches(0.4),
    [[("폭발적인 상업적 성공 — 역대 가장 빠르게 성장한 소프트웨어 기업 중 하나:", 15, NAVY, True)]])
# valuation timeline chips
tl = [
    ("$9.9B","2025년 6월 · ARR 5억 달러 이상", TEAL),
    ("$29.3B","2025년 11월 · 시리즈 D, 23억 달러 조달", GREEN),
    ("$2B ARR","2026년 2월 · 약 3년 만에 달성", AMBER),
    ("$60B","2026년 6월 · SpaceX가 인수", CORAL),
]
tmargin=Inches(0.7); tgap=Inches(0.3); tcw=(SW-2*tmargin-3*tgap)/4; tch=Inches(1.35); ty=Inches(3.75)
for i,(big,sub,col) in enumerate(tl):
    x=tmargin+i*(tcw+tgap)
    rect(s,x,ty,tcw,tch,CARD,line=LINE,line_w=1,shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s,x,ty,tcw,Inches(0.12),col)
    txt(s,x+Inches(0.15),ty+Inches(0.22),tcw-Inches(0.3),tch-Inches(0.3),
        [[(big,26,col,True)],[(sub,11,GRAY,False)]],align=PP_ALIGN.CENTER,
        space_after=6,line_spacing=1.05)
    if i<3:
        rect(s,x+tcw+Inches(0.02),ty+tch/2-Inches(0.12),Inches(0.26),Inches(0.24),
             col,shape=MSO_SHAPE.RIGHT_ARROW)
note(s, "SpaceX가 Cursor(Anysphere)를 $60B에 인수해 xAI 산하로 편입 — 'AI 네이티브' IDE가 하나의 거대 산업이 되었습니다.", y=Inches(5.55))
footer(s)

# =========================================================
# SLIDE 6 — Agentic AI -> autonomous coding
# =========================================================
s = slide(); bg(s); header(s, 6, "에이전틱 AI — Claude Code와 Codex가 도구를 직접 다룬다", "2024~2025년  ·  자율 코딩(autonomous coding)의 시대")
flow(s, [("VS Code / 터미널","호스트"),
         ("에이전틱 AI","Claude Code · Codex"),
         ("자율 코딩","계획 · 수정 · 테스트 · 배포")],
     y=Inches(2.5), h=Inches(1.9), colors=[NAVY, TEAL, CORAL])
bullets(s, Inches(0.7), Inches(4.7), Inches(11.8), [
    ("AI가 더 이상 제안만 하지 않습니다 — 파일을 읽고, 명령을 실행하고, 코드를 고치고, 결과를 검증합니다.", ""),
    ("하나의 지시가 여러 자율 단계로 확장됩니다: 에이전트가 전체 루프를 스스로 이끕니다.", ""),
], size=15, gap=12)
note(s, "제안(suggest)을 넘어 실행(act)으로 — 에이전트가 파일을 읽고, 명령을 실행하고, 코드를 고치는 '자율 코딩' 시대.", y=Inches(5.9))
footer(s)

# =========================================================
# SLIDE 7 — AntiGravity (Google)
# =========================================================
s = slide(); bg(s); header(s, 7, "Antigravity — 구글의 브랜드 에이전틱 IDE", "2025년  ·  플랫폼 대기업들의 진입")
flow(s, [("VS Code 계보","익숙한 기반"),
         ("Google Antigravity","브랜딩 & 커스터마이즈"),
         ("에이전트 우선 IDE","구글 모델에 최적화")],
     y=Inches(2.35), h=Inches(1.75), colors=[NAVY, GREEN, TEAL])
bullets(s, Inches(0.7), Inches(4.45), Inches(11.8), [
    ("이제 대형 벤더들도 VS Code 계보 위에 자신만의 관점을 담은 브랜드 AI IDE를 내놓습니다.", ""),
    ("패턴은 정해졌습니다:  VS Code를 템플릿으로 삼아, 특정 사용자층에 맞게 커스터마이즈한다.", ""),
], size=15, gap=11)
note(s, "핵심 흐름: VS Code를 '템플릿'으로 삼아 특정 목적에 맞게 브랜딩·커스터마이즈한다 — 다음은 '연구자'입니다.", y=Inches(5.75))
footer(s)

# =========================================================
# SLIDE 8 — GHBIO: IDE for Omics researcher
# =========================================================
s = slide(); bg(s, NAVY)
kicker_bar(s)
rect(s, 0, Inches(2.0), SW, Inches(3.5), NAVY2)
txt(s, Inches(0.8), Inches(0.9), Inches(11.8), Inches(0.5),
    [[("다음 단계", 15, TEAL, True)]])
txt(s, Inches(0.8), Inches(1.4), Inches(11.8), Inches(1.6),
    [[("GHBio가 오믹스 연구자를 위한 IDE를 엽니다", 34, WHITE, True)],
     [("오믹스 연구자를 위한 AI IDE — VS Code를 템플릿으로", 18, RGBColor(0xBF,0xD6,0xE0), False)]],
    space_after=8)
flow(s, [("VS Code / code-server","검증된 템플릿"),
         ("+ 오믹스 파이프라인","scRNA-seq, 그리고 그 너머"),
         ("GHBio Co-Scientist","생물학자를 위한 IDE")],
     y=Inches(3.5), h=Inches(1.55), colors=[TEAL, GREEN, AMBER])
txt(s, Inches(0.8), Inches(5.5), Inches(11.8), Inches(1.2),
    [[("Cursor·Antigravity와 같은 발상 — 다만 대상은 소프트웨어 엔지니어가 아니라 ", 15, WHITE, False),
      ("실험실(wet-lab)의 생물학자입니다.", 15, TEAL, True)]],
    anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.15)
footer(s)

# =========================================================
# SLIDE 9 — GHBio Co-Scientist overview
# =========================================================
s = slide(); bg(s); header(s, 9, "GHBio Co-Scientist — 오믹스를 위한 작업대", "제품 소개  ·  생물학을 위한 PlatformIO 스타일")
bullets(s, Inches(0.7), Inches(1.7), Inches(6.2), [
    ("VS Code / code-server 확장 프로그램", "— 브라우저에서 실행, 설치할 것이 없습니다."),
    ("PlatformIO 스타일의 생명정보학 작업대", "프로그래머가 아닌 생물학자를 위한."),
    ("트리 뷰 + 신뢰할 수 있는 AI 분석 패널", "분석 전 과정을 처음부터 끝까지 안내."),
    ("이중 언어 UI (한국어 우선)", "코딩 경험 없는 실험실 연구자를 위해 설계."),
], size=14, gap=13)
# right-half stacked cards
rc = [
    ("브라우저에서 실행","ghbiocosci.iotok.org의 code-server — 탭만 열면 바로 분석 시작."),
    ("코딩 불필요","파이프라인 단계를 클릭하면 확장이 대신 도구를 실행합니다."),
    ("AI Co-Scientist 내장","결과에 대한 생물학적 해석을 원클릭으로."),
]
rx = Inches(7.15); rw = Inches(5.5); ry = Inches(1.75); rh = Inches(1.35); rgap = Inches(0.22)
for i,(t,b) in enumerate(rc):
    yy = ry + i*(rh+rgap)
    rect(s, rx, yy, rw, rh, CARD, line=LINE, line_w=1, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s, rx, yy, Inches(0.12), rh, [TEAL, GREEN, AMBER][i])
    txt(s, rx+Inches(0.3), yy+Inches(0.16), rw-Inches(0.5), rh-Inches(0.3),
        [[(t, 15, NAVY, True)], [(b, 12, GRAY, False)]], space_after=5, line_spacing=1.06)
footer(s)

# =========================================================
# SLIDE 10 — Who it's for / the problem
# =========================================================
s = slide(); bg(s); header(s, 10, "이 제품이 푸는 문제", "사용자  ·  Copilot 없는 한국 생물학자")
cards(s, [
    ("간극","현대 오믹스 분석은 명령줄(command line)에서 이뤄집니다 — 정렬, Python, 클러스터링. 대부분의 생물학자는 코딩을 하지 않습니다."),
    ("장벽","STAR·레퍼런스·Scanpy·GPU 서버 세팅은 생물학을 시작하기도 전에 며칠간의 DevOps 작업입니다."),
    ("해답","Co-Scientist는 이 모든 것을 '클릭하면 실행되는' 단계와, 결과를 한국어로 설명하는 AI 뒤에 숨깁니다."),
], y=Inches(1.75), h=Inches(2.5), cols=3, top_colors=[CORAL, AMBER, GREEN])
bullets(s, Inches(0.7), Inches(4.6), Inches(11.8), [
    ("이중 언어 UI 문구(한국어 우선), AI 패널은 기본적으로 한국어로 답변합니다.", ""),
    ("대용량 입력(FASTQ, 약 30~40 GB GRCh38 인덱스)은 한 번만 세팅해 재사용 — 반복 다운로드 없음.", ""),
], size=14, gap=11)
footer(s)

# =========================================================
# SLIDE 11 — Architecture: data-driven registry
# =========================================================
s = slide(); bg(s); header(s, 11, "아키텍처 — 데이터 주도(data-driven) 레지스트리", "설계 목표  ·  파일만 넣으면 새 도메인 추가")
txt(s, Inches(0.7), Inches(1.6), Inches(11.8), Inches(0.5),
    [[("새로운 분석 도메인은 JSON 매니페스트로 추가됩니다 — ", 15, INK, False),
      ("TypeScript 코드 수정 없이.", 15, CORAL, True)]])
flow(s, [("module.json","도메인 정체성 + 라이브러리 + AI 설정"),
         ("pipeline.json","파이프라인별 순서가 있는 단계"),
         ("트리 + AI 패널","매니페스트로부터 런타임에 렌더링")],
     y=Inches(2.3), h=Inches(1.7), colors=[NAVY, TEAL, GREEN])
cards(s, [
    ("modules/<id>/module.json","정체성, 설치 가능한 라이브러리(presentPath 존재 확인 포함), 그리고 AI 설정."),
    ("modules/<id>/pipelines/<pid>/pipeline.json","파이프라인의 순서가 있는 단계(task / ai 종류)."),
    ("_ 접두사 = 무시됨","_template / _wip-* 같은 디렉터리는 로더가 건너뜁니다."),
], y=Inches(4.3), h=Inches(2.1), cols=3, top_colors=[TEAL, GREEN, AMBER])
footer(s)

# =========================================================
# SLIDE 12 — Module registry
# =========================================================
s = slide(); bg(s); header(s, 12, "모듈 레지스트리 (src/modules.ts)", "하나의 모듈 = 하나의 분석 도메인")
bullets(s, Inches(0.7), Inches(1.7), Inches(6.1), [
    ("모듈은 하나의 완결된 분석 도메인", "— 오늘은 scRNA-seq, 내일은 단백질 모델링."),
    ("도메인 고유의 모든 것은 JSON", "하드코딩이 아니라 런타임에 읽습니다."),
    ("자신의 라이브러리를 선언", "설치 가능한 도구 + 존재 확인(presentPath)."),
    ("자신의 AI 설정을 선언", "시스템 프롬프트, 프리셋 프롬프트, 결과 파일 컨텍스트."),
], size=14, gap=12)
# code-ish card
rect(s, Inches(7.05), Inches(1.7), Inches(5.6), Inches(4.4), NAVY,
     shape=MSO_SHAPE.ROUNDED_RECTANGLE)
txt(s, Inches(7.05), Inches(1.75), Inches(5.6), Inches(0.4),
    [[("modules/scrna-seq/module.json", 12, TEAL, True)]], align=PP_ALIGN.CENTER)
code_lines = [
    '"id": "scrna-seq",',
    '"name": "Single-cell RNA-seq",',
    '"libraries": [',
    '   Python / Scanpy   (~/ghbio-venv)',
    '   STAR / STARsolo   (~/bin/STAR)',
    '   GRCh38 인덱스     (~30 GB)',
    '],',
    '"ai": { system, prompts,',
    '        context, readyFile }',
]
txt(s, Inches(7.35), Inches(2.25), Inches(5.1), Inches(3.7),
    [[(l, 13, RGBColor(0xCF,0xE8,0xE4), False)] for l in code_lines],
    space_after=6, line_spacing=1.1)
footer(s)

# =========================================================
# SLIDE 13 — Three tree views
# =========================================================
s = slide(); bg(s); header(s, 13, "세 개의 트리 뷰가 앱 전체를 이끈다", "파이프라인 · 프로젝트 · 라이브러리")
cards(s, [
    ("🧪  파이프라인 (Pipelines)","분석을 골라 단계를 순서대로 실행 — 터미널에 ▶ / ✅ 배너가 표시됩니다."),
    ("📁  프로젝트 (Projects)","파이프라인마다 전용 프로젝트 폴더; 결과는 탐색 가능한 일급 파일입니다."),
    ("📦  라이브러리 (Libraries)","설치 가능한 도구(Scanpy, STAR, GRCh38 인덱스)와 실시간 '설치됨?' 확인."),
], y=Inches(1.9), h=Inches(3.3), cols=3, top_colors=[TEAL, GREEN, AMBER])
note(s, "세 개의 트리 뷰가 UI의 전부입니다 — 파이프라인 실행, 결과 탐색, 도구 설치가 클릭 한 번으로.", y=Inches(5.5))
footer(s)

# =========================================================
# SLIDE 14 — Pipeline engine
# =========================================================
s = slide(); bg(s); header(s, 14, "파이프라인 엔진 (src/pipeline.ts)", "핵심 '딥 모듈(deep module)'")
bullets(s, Inches(0.7), Inches(1.7), Inches(6.1), [
    ("풍부한 엔진 위의 좁은 인터페이스", "runPipeline · stageStatus · isStageDone · ensureProject."),
    ("어려운 부분을 숨깁니다", "단계 순서, 프로젝트 구성, 멱등성(idempotency), AI 인계."),
    ("내부에 UI 관심사 없음", "덕분에 향후 헤드리스 'AI-as-a-webservice'도 지원할 수 있습니다."),
], size=14, gap=13)
cards(s, [
    ('kind: "task"','셸 명령을 VS Code Task로 터미널에서 실행 — ▶ / ✅ / → 다음 배너와 함께.'),
    ('kind: "ai"','결과를 해석하기 위해 AI 분석 패널을 엽니다.'),
    ('produces: [...]','선언된 산출물 → 해당 단계를 ✓ 완료로 자동 표시(멱등적 재실행).'),
], y=Inches(1.75), h=Inches(4.3), cols=1, top_colors=[TEAL, AMBER, GREEN])
footer(s)

# =========================================================
# SLIDE 15 — GHBIO_RESULTS
# =========================================================
s = slide(); bg(s); header(s, 15, "GHBIO_RESULTS — 결과가 곧 프로젝트 파일", "핵심을 떠받치는 규칙")
txt(s, Inches(0.7), Inches(1.65), Inches(11.8), Inches(0.5),
    [[("모든 단계는 ", 15, INK, False),
      ("GHBIO_RESULTS", 15, CORAL, True),
      (" 를 해당 파이프라인 전용 결과 폴더로 가리킨 채 실행됩니다:", 15, INK, False)]])
rect(s, Inches(0.7), Inches(2.25), Inches(11.9), Inches(0.7), NAVY,
     shape=MSO_SHAPE.ROUNDED_RECTANGLE)
txt(s, Inches(0.9), Inches(2.25), Inches(11.5), Inches(0.7),
    [[("~/ghbio-workspace/projects/<pipelineId>/results/", 17, TEAL, True)]],
    anchor=MSO_ANCHOR.MIDDLE)
bullets(s, Inches(0.7), Inches(3.35), Inches(11.8), [
    ("파이프라인이 실행하는 모든 스크립트는 반드시 GHBIO_RESULTS를 읽어야 합니다 — 경로를 하드코딩하면 엉뚱한 프로젝트에 씁니다.", ""),
    ("결과는 일급 파일이 되어 프로젝트 뷰에 나타나고 AI 패널에 입력됩니다.", ""),
    ("대용량 공유 입력(FASTQ, GRCh38 인덱스, 약 40 GB)은 ~/ghbio-tutorial/ 아래 두고 멱등적으로 재사용합니다.", ""),
], size=14, gap=13)
note(s, "결과물이 곧 프로젝트 파일입니다 — 이 규칙 덕분에 AI 분석과 리포트 단계가 항상 올바른 출력을 찾습니다.", y=Inches(5.85))
footer(s)

# =========================================================
# SLIDE 16 — scRNA-seq module libraries
# =========================================================
s = slide(); bg(s); header(s, 16, "scRNA-seq 모듈 — 설치하는 도구들", "단일세포 RNA-seq  ·  aarch64용으로 빌드")
cards(s, [
    ("Python / Scanpy","~/ghbio-venv — scanpy, anndata, leiden, umap. 원클릭 환경 세팅."),
    ("STAR / STARsolo","~/bin/STAR — 소스에서 컴파일한 정렬기(ARM에는 Cell Ranger 없음)."),
    ("GRCh38 레퍼런스","~/ghbio-tutorial/ref/star_index — 약 30 GB 인덱스, 한 번 빌드 후 재사용."),
], y=Inches(1.85), h=Inches(2.6), cols=3, top_colors=[GREEN, TEAL, AMBER])
bullets(s, Inches(0.7), Inches(4.75), Inches(11.8), [
    ("각 라이브러리는 presentPath를 선언 — 라이브러리 뷰가 설치 여부를 실시간으로 보여줍니다.", ""),
    ("두 데이터셋이 파이프라인으로 제공: PBMC(10x)와 교모세포종(Glioblastoma, 10x); 여기에 '내 매트릭스 가져오기'까지.", ""),
], size=14, gap=11)
footer(s)

# =========================================================
# SLIDE 17 — pipeline stages
# =========================================================
s = slide(); bg(s); header(s, 17, "하나의 파이프라인: FASTQ → 매트릭스 → 클러스터링 → AI", "scRNA-seq PBMC / 교모세포종  ·  8단계")
stages = [
    ("0 · 셋업","Python 환경"),
    ("1 · 다운로드","FASTQ 데이터"),
    ("2a · STAR 빌드","정렬기 컴파일"),
    ("2b · 레퍼런스","GRCh38 인덱스"),
    ("2c · STARsolo","정렬 → 매트릭스"),
    ("3 · Scanpy","QC & 클러스터링"),
    ("4 · AI","가설 도출"),
    ("5 · 리포트","PDF"),
]
# two rows of 4
margin=Inches(0.6); gap=Inches(0.3); cw=(SW-2*margin-3*gap)/4; ch=Inches(1.25)
pal=[NAVY,NAVY2,TEAL,TEAL,GREEN,GREEN,AMBER,CORAL]
for i,(a,b) in enumerate(stages):
    r=i//4; c=i%4
    x=margin+c*(cw+gap); y=Inches(1.9)+r*(ch+Inches(0.55))
    rect(s,x,y,cw,ch,pal[i],shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(s,x+Inches(0.1),y,cw-Inches(0.2),ch,
        [[(a,15,WHITE,True)],[(b,11,RGBColor(0xDF,0xEC,0xF0),False)]],
        align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,space_after=3,line_spacing=1.0)
    if c<3:
        rect(s,x+cw+Inches(0.02),y+ch/2-Inches(0.14),Inches(0.26),Inches(0.28),
             AMBER,shape=MSO_SHAPE.RIGHT_ARROW)
note(s, "'내 매트릭스 가져오기' 파이프라인은 정렬 단계를 건너뛰고 사용자의 count matrix에서 바로 시작합니다.", y=Inches(5.65))
footer(s)

# =========================================================
# SLIDE 18 — AI Analysis panel
# =========================================================
s = slide(); bg(s); header(s, 18, "AI 분석 패널 — 설계부터 신뢰할 수 있게", "src/ai/panel.ts  ·  단발 요청, 에이전트 루프 없음")
bullets(s, Inches(0.7), Inches(1.7), Inches(6.2), [
    ("하나의 스트리밍 요청 → 답변 렌더링.", "에이전트·도구·파일 수정 없음."),
    ("Stop 버튼으로 취소 가능.", "이 '신뢰성 확보'가 핵심 목적입니다."),
    ("결과 CSV를 프롬프트에 읽어 들임", "markers_by_cluster.csv, celltype_draft.csv."),
    ("기본적으로 한국어로 답변.", "신중한 단일세포 분석가 페르소나."),
], size=14, gap=12)
cards(s, [
    ("프로바이더","anthropic · groq · openrouter · deepseek (OpenAI 호환). API 키는 저장소 밖에 보관."),
    ("프리셋 vs. 자유형","프리셋은 결과 파일이 필요; 자유형 질문은 언제든 답변."),
    ("의도적으로 단순하게","여기에는 도구·에이전트 루프가 없습니다 — 그래서 안정적입니다."),
], y=Inches(1.75), h=Inches(4.3), cols=1, top_colors=[TEAL, GREEN, AMBER])
footer(s)

# =========================================================
# SLIDE 19 — Preset prompts & save to report
# =========================================================
s = slide(); bg(s); header(s, 19, "프리셋 프롬프트 & 리포트로 저장", "마커에서 가설까지, 클릭 한 번에")
cards(s, [
    ("클러스터 해석","클러스터별 세포 정체성·활성화 상태를 신뢰도와 이중체(doublet) 표시와 함께."),
    ("주석 다듬기","표준 PBMC 마커 → 깔끔한 클러스터별 세포유형 표."),
    ("경로(Pathway)","마커를 생물학적 프로그램으로 묶고, 검증할 유전자 세트를 제안."),
    ("검증 가능한 가설","근거·예측·검증 방법이 담긴 가설 3~5개."),
    ("질문 목록","우선순위가 매겨진 후속 분석 8~10개, 각각 한 줄 방법 포함."),
    ("📘 쉬운 리포트","일상적 비유로 가득한, 고등학생도 이해할 리포트(한국어)."),
], y=Inches(1.75), h=Inches(3.1), cols=3, top_colors=[TEAL,GREEN,AMBER,CORAL,NAVY2,TEAL])
note(s, "'리포트로 저장' → step4_ai_report.md / step4_ai_report_easy.md 를 결과 폴더에 저장하고, 05_make_report.sh 가 PDF에 선택적으로 포함합니다.", y=Inches(5.35))
footer(s)

# =========================================================
# SLIDE 20 — Deployment
# =========================================================
s = slide(); bg(s); header(s, 20, "배포 방식 — ARM 위의 code-server", "배포  ·  브라우저 네이티브, 자체 호스팅")
cards(s, [
    ("code-server","ghbiocosci.iotok.org에서 브라우저로 실행 — VS Code, 로컬 설치 불필요."),
    ("aarch64 서버","ARM 서버: STAR는 소스에서 빌드, Cell Ranger 없음. Python venv는 ~/ghbio-venv."),
    ("build.sh 배포","esbuild 번들 → .vsix 수동 zip → 설치 → ghbio-code 서비스 재시작."),
], y=Inches(1.85), h=Inches(2.6), cols=3, top_colors=[TEAL, GREEN, AMBER])
bullets(s, Inches(0.7), Inches(4.75), Inches(11.8), [
    ("code-server 설정은 내장 Copilot/채팅을 끄고, 대용량 *.fastq / *.bam 파일을 탐색기에서 제외합니다.", ""),
    ("배포마다 버전 상향 + 브라우저 강력 새로고침 필요 — AI 패널이 유일한 접점입니다.", ""),
], size=14, gap=11)
footer(s)

# =========================================================
# SLIDE 21 — Conclusion
# =========================================================
s = slide(); bg(s, NAVY); kicker_bar(s)
rect(s, 0, Inches(1.7), SW, Inches(0.06), TEAL)
txt(s, Inches(0.8), Inches(0.8), Inches(11.8), Inches(0.9),
    [[("결론", 15, TEAL, True)],
     [("IDE가 실험실(wet lab)을 만나다", 34, WHITE, True)]], space_after=6)
cards2 = [
    ("같은 계보","VS Code → Copilot → Cursor → 에이전틱 IDE. Co-Scientist는 그 검증된 길 위에 섭니다."),
    ("새로운 대상","소프트웨어 엔지니어가 아니라 생물학자. 클릭 실행 파이프라인과, 결과를 설명하는 AI."),
    ("데이터 주도","module.json + pipeline.json만 넣으면 완전히 새로운 오믹스 도메인이 나타납니다. 코딩 없이."),
]
margin=Inches(0.8); gap=Inches(0.4); cw=(SW-2*margin-2*gap)/3
for i,(t,b) in enumerate(cards2):
    x=margin+i*(cw+gap); y=Inches(2.1)
    rect(s,x,y,cw,Inches(2.5),NAVY2,shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s,x,y,cw,Inches(0.12),[TEAL,GREEN,AMBER][i])
    txt(s,x+Inches(0.25),y+Inches(0.3),cw-Inches(0.5),Inches(2.0),
        [[(t,17,WHITE,True)],[(b,13,RGBColor(0xC7,0xDA,0xE2),False)]],
        space_after=8,line_spacing=1.12)
txt(s, Inches(0.8), Inches(5.1), Inches(11.8), Inches(1.4),
    [[("오믹스 연구자를 위한 AI IDE — ", 20, WHITE, True),
      ("IDE for Omics Researcher", 20, TEAL, True)],
     [("자율 코딩 혁명을, 생물학적 질문을 던지는 사람들에게로.", 14, RGBColor(0xBF,0xD6,0xE0), False)]],
    space_after=10, line_spacing=1.2)
footer(s)

# =========================================================
# SLIDE 22 — About Us
# =========================================================
s = slide(); bg(s); header(s, 22, "회사 소개 — 지에이치바이오(주) / GHBio", "회사  ·  ghbio.co.kr")
bullets(s, Inches(0.7), Inches(1.7), Inches(6.3), [
    ("지에이치바이오(주) (GHBio)", "신약개발 연구를 위한 혁신적 마우스 모델 전문기업."),
    ("전임상 유효성 & 안전성", "in-vivo 유효성, 비-GLP 독성·약리, 약동학(PK), 생물발광 이미징."),
    ("유전자 편집 & 인간화 모델", "대사질환, 바이러스 감염, 면역결핍, 인간화 간(humanized-liver) 모델."),
    ("면역관문 인간화 마우스", "PD-1, CTLA-4, TIGIT, TIM-3, 4-1BB, CD3E — 항암 면역치료 연구용."),
], size=13.5, gap=12)
# contact card
rect(s, Inches(7.35), Inches(1.7), Inches(5.3), Inches(4.4), CARD, line=LINE, line_w=1,
     shape=MSO_SHAPE.ROUNDED_RECTANGLE)
rect(s, Inches(7.35), Inches(1.7), Inches(5.3), Inches(0.12), TEAL)
txt(s, Inches(7.65), Inches(2.0), Inches(4.7), Inches(4.0),
    [[("연락처", 16, NAVY, True)],
     [("대표", 12, TEAL, True),("   유경원 (Yoo Kyung-won)", 13, INK, False)],
     [("주소", 12, TEAL, True),("   대전 유성구 테크노4로 17, D217", 13, INK, False)],
     [("", 12, TEAL, True),("   대한민국 대전광역시", 12, GRAY, False)],
     [("전화", 12, TEAL, True),("   042-716-2177", 13, INK, False)],
     [("이메일", 12, TEAL, True),("   ghbio@gh-bio.com", 13, INK, False)],
     [("웹", 12, TEAL, True),("   ghbio.co.kr", 13, INK, False)],
     [("사업자등록번호", 12, TEAL, True),("   318-81-03789", 13, INK, False)],
    ], space_after=11, line_spacing=1.1)
note(s, "신약개발용 마우스 모델과 전임상 유효성·안전성 평가 서비스 전문기업 — 그 데이터 과학 역량이 GHBio Co-Scientist로 이어집니다.", y=Inches(6.2))
footer(s)

out = "/home/jit/ghbio-coscientist/IDE_for_Omics_Researcher_KO.pptx"
prs.save(out)
print("SAVED", out, "slides:", len(prs.slides._sldIdLst))
